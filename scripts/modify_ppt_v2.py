#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
修改湖蓝PPT Slide 4 + Slide 5
- Slide 4: 拆分痛点/解决思路，删除重复内容
- Slide 5: 优化国内外研究现状结构
"""

import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor

src = '湖蓝 基于生成对抗网络的工业异常检测图像样本增广研究与实现.pptx'
dst = '湖蓝 基于生成对抗网络的工业异常检测图像样本增广研究与实现-修改版v2.pptx'

prs = Presentation(src)

# ============================================================
# SLIDE 4: 拆分 "痛点" 与 "解决思路"
# ============================================================
slide4 = prs.slides[3]

# --- Step 1: 删除重复的"深度学习模型训练困难" ---
to_delete = []
for i, shape in enumerate(slide4.shapes):
    if shape.has_text_frame:
        t = shape.text_frame.text
        # 重复的标题位置更靠下 (top > 3500000)
        if '深度学习模型训练困难' in t and shape.top > 3500000 and shape.width > 10000000:
            to_delete.append(i)
        # 以及紧跟它的描述框
        if '模型长期看到大量正常样本' in t and shape.top > 4000000:
            to_delete.append(i)

for idx in sorted(to_delete, reverse=True):
    sp = slide4.shapes[idx]._element
    sp.getparent().remove(sp)

print(f'Slide 4: 删除了 {len(to_delete)} 个重复形状')

# --- Step 2: 修改标题，突出痛点/方案双栏结构 ---
for shape in slide4.shapes:
    if not shape.has_text_frame:
        continue
    t = shape.text_frame.text

    # 主标题修改
    if t.strip() == '工业生产中的现实痛点与解决思路':
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '工业生产中的现实痛点与解决思路' in run.text:
                    run.text = '现实痛点分析 → 解决思路：基于GAN的样本增广'

    # 把"技术方案：GAN样本增广" 改为更清晰的表述
    if t.strip() == '技术方案：GAN样本增广':
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '技术方案：GAN样本增广' in run.text:
                    run.text = '解决思路：基于GAN的伪异常样本生成'

    # 底部 Focus-StyleGAN 文本改为总结性表述，不抢占后面页面的内容
    if '本文提出Focus-StyleGAN系统' in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '本文提出Focus-StyleGAN系统' in run.text:
                    run.text = run.text.replace(
                        '本文提出Focus-StyleGAN系统，融合FocusGAN缺陷聚焦能力与StyleGAN风格控制优势，解决工业异常检测样本不平衡问题。',
                        '核心思路：利用GAN的生成能力，在仅依赖正常样本的条件下，合成形态逼真的伪异常图像，从根本上缓解工业检测中的数据不平衡问题。'
                    )
                    break

    # 优化"理论意义与应用价值"描述
    if '双分支解耦架构提供可控异常合成新范式' in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '双分支解耦架构提供可控异常合成新范式' in run.text:
                    run.text = run.text.replace(
                        '双分支解耦架构提供可控异常合成新范式；可用于现有工业视觉检测系统，大大减少对稀缺异常样本的依赖，推动智能制造质量控制全流程智能化。',
                        '理论价值：双分支解耦架构为可控异常合成提供新范式。应用价值：可接入现有工业视觉检测系统，减少对稀缺异常样本的依赖，降低企业数据采集与标注成本。'
                    )
                    break

print('Slide 4: 标题和关键文本修改完成')

# ============================================================
# SLIDE 5: 国内外研究现状 - 优化结构
# ============================================================
slide5 = prs.slides[4]

# --- 修正研究空白表述 ---
for shape in slide5.shapes:
    if not shape.has_text_frame:
        continue
    t = shape.text_frame.text

    if '研究空白' in t and '均出现性能衰减' in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '均出现性能衰减' in run.text:
                    run.text = run.text.replace(
                        '现有方法在处理极端样本不平衡时均出现性能衰减，研究思路已转向从数据源头主动生成高质量伪异常样本。',
                        '现有工作在极端样本不平衡（正常:异常>100:1）时性能普遍衰减。近年来研究重心正从改进检测算法转向从数据源头生成高质量伪异常样本——这正是本课题的切入点。'
                    )
                    break

    # 修正国外-AnoGAN条目：补充与本课题关系
    if 'AnoGAN首次将GAN引入异常检测任务' in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if 'AnoGAN首次将GAN引入异常检测任务' in run.text:
                    run.text = run.text.replace(
                        'AnoGAN首次将GAN引入异常检测任务',
                        'AnoGAN（2018）首次将GAN用于异常检测，但仅做检测未做增广'
                    )
                    break

    # 修正国外-OCGAN条目
    if 'OCGAN用约束潜在表示做新颖性检测' in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if 'OCGAN用约束潜在表示做新颖性检测' in run.text:
                    run.text = run.text.replace(
                        'OCGAN用约束潜在表示做新颖性检测',
                        'OCGAN（CVPR 2019）约束潜在空间做单类检测，启发了本课题的潜在空间调制思路'
                    )
                    break

    # 修正国内-尹高科条目
    if '尹高科：改进GAN无监督工业图像异常检测' in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '尹高科：改进GAN无监督工业图像异常检测' in run.text:
                    run.text = run.text.replace(
                        '尹高科：改进GAN无监督工业图像异常检测，减少标注依赖',
                        '尹高科：改进GAN+数据增强的无监督检测，将增广与检测结合'
                    )
                    break

    # 补充CBAM相关
    if '王小雅：CBAM与Cascade RPN融合的形变识别方法' in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '王小雅：CBAM与Cascade RPN融合的形变识别方法' in run.text:
                    run.text = run.text.replace(
                        '王小雅：CBAM与Cascade RPN融合的形变识别方法',
                        '王小雅：CBAM注意力+Cascade RPN融合，验证了CBAM在工业视觉中的有效性'
                    )
                    break

print('Slide 5: 国内外研究现状优化完成')

# ============================================================
# SAVE
# ============================================================
prs.save(dst)
print(f'\n已保存至: {dst}')
