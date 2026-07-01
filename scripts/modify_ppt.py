#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""修改湖蓝PPT的内容"""

import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from copy import deepcopy

src = '湖蓝 基于生成对抗网络的工业异常检测图像样本增广研究与实现.pptx'
dst = '湖蓝 基于生成对抗网络的工业异常检测图像样本增广研究与实现-修改版.pptx'

prs = Presentation(src)

# ============================================================
# SLIDE 4 (index 3): 背景页 - 删除重复文本框 + 修改过早出现的模型名
# ============================================================
slide4 = prs.slides[3]
to_delete = []
for i, shape in enumerate(slide4.shapes):
    text = shape.text_frame.text if shape.has_text_frame else ''
    # 重复的"深度学习模型训练困难"——位置靠下的那个
    if '深度学习模型训练困难' in text and shape.top > 3500000:
        to_delete.append(i)
    # "本文提出Focus-StyleGAN"出现太早，改为通用的解决思路
    if '本文提出Focus-StyleGAN系统' in text:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '本文提出Focus-StyleGAN系统' in run.text:
                    run.text = run.text.replace(
                        '本文提出Focus-StyleGAN系统，融合FocusGAN缺陷聚焦能力与StyleGAN风格控制优势，解决工业异常检测样本不平衡问题。',
                        '针对上述痛点，本课题研究在仅依赖正常样本的条件下，利用生成对抗网络实现高保真伪异常图像生成，解决工业异常检测样本不平衡问题。'
                    )
                    break

for idx in sorted(to_delete, reverse=True):
    sp = slide4.shapes[idx]._element
    sp.getparent().remove(sp)

print('Slide 4: 删除重复文本框，修正过早出现的模型名')

# ============================================================
# SLIDE 5 (index 4): 国内外研究现状 - 修正研究空白表述
# ============================================================
slide5 = prs.slides[4]
for shape in slide5.shapes:
    if shape.has_text_frame and '研究空白' in shape.text_frame.text:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '均出现性能衰减' in run.text:
                    run.text = run.text.replace(
                        '现有方法在处理极端样本不平衡时均出现性能衰减，研究思路已转向从数据源头主动生成高质量伪异常样本。',
                        '现有方法在极端样本不平衡场景（正常:异常>100:1）下性能普遍衰减，增广式方法尚处于探索阶段，本课题聚焦此空白展开研究。'
                    )
                    break

print('Slide 5: 修正研究空白表述')

# ============================================================
# SLIDE 7 (index 6): 研究目标 - 修正底部注释 + 评估描述
# ============================================================
slide7 = prs.slides[6]
for shape in slide7.shapes:
    if shape.has_text_frame and '注：当前实验验证' in shape.text_frame.text:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '注：当前实验验证' in run.text:
                    run.text = run.text.replace(
                        '注：当前实验验证主要聚焦于第一阶段静态批量增广，完整的闭环动态优化将在后续工作中实现。',
                        '注：本课题已完成完整的静态批量增广、质量评估与检测验证闭环流程，并构建了Web演示系统，闭环动态优化作为后续扩展方向。'
                    )
                    break

for shape in slide7.shapes:
    if shape.has_text_frame and '建立图像质量(FID/IS/LPIPS)与检测性能' in shape.text_frame.text:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '建立图像质量(FID/IS/LPIPS)与检测性能' in run.text:
                    run.text = run.text.replace(
                        '建立图像质量(FID/IS/LPIPS)与检测性能(Pixel-AUC/PRO-AUC)双重评估体系',
                        '建立图像质量(FID/IS/LPIPS/PPS)与检测性能(Pixel-AUC/PRO-AUC)双重评估，并构建Web演示系统'
                    )
                    break

print('Slide 7: 修正底部注释和评估描述')

# ============================================================
# SLIDE 8 (index 7): 主要研究内容 - 修正科学问题 + 优化四项描述
# ============================================================
slide8 = prs.slides[7]
for shape in slide8.shapes:
    t = shape.text_frame.text if shape.has_text_frame else ''

    if '核心科学问题为' in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '核心科学问题为' in run.text:
                    run.text = run.text.replace(
                        '核心科学问题为：工业场景下高质量可控伪异常样本的生成。',
                        '核心科学问题：如何在仅依赖正常样本的条件下，实现缺陷类型可控、背景结构保持的高保真伪异常图像生成？'
                    )
                    break

    if '缺陷聚焦分支生成缺陷区域特征' in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '缺陷聚焦分支生成缺陷区域特征' in run.text:
                    run.text = run.text.replace(
                        '缺陷聚焦分支生成缺陷区域特征，背景保持分支维护正常区域，注意力引导融合模块实现自然结合',
                        '缺陷聚焦分支生成缺陷纹理（集成CBAM注意力），背景保持分支维护产品结构，融合模块自适应结合双路输出'
                    )
                    break

    if '三个尺度子判别器' in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '三个尺度子判别器' in run.text:
                    run.text = run.text.replace(
                        '三个尺度子判别器（原图/半分辨率/四分之一分辨率），PatchGAN结构，CBAM集成增强微缺陷敏感度',
                        '三尺度子判别器并行（原图/半分辨率/1/4分辨率），每层嵌入CBAM通道+空间注意力，增强微缺陷判别力'
                    )
                    break

    if 'WGAN-GP对抗损失+感知损失+重构损失+LPIPS损失' in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if 'WGAN-GP对抗损失+感知损失+重构损失+LPIPS损失' in run.text:
                    run.text = run.text.replace(
                        'WGAN-GP对抗损失+感知损失+重构损失+LPIPS损失，从像素/感知/语义多层面约束生成过程',
                        'WGAN-GP对抗 + VGG19感知 + L1重构 + LPIPS感知相似度，四损失联合约束生成过程'
                    )
                    break

    if '图像质量指标(FID/IS/LPIPS/PPS)与下游检测性能' in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '图像质量指标(FID/IS/LPIPS/PPS)与下游检测性能' in run.text:
                    run.text = run.text.replace(
                        '图像质量指标(FID/IS/LPIPS/PPS)与下游检测性能(Pixel-AUC/PRO-AUC)双重验证增广效果',
                        '图像质量(FID/IS/LPIPS/PPS) + 检测性能(Pixel-AUC/PRO-AUC) + Web交互，三维度验证增广效果'
                    )
                    break

print('Slide 8: 修正科学问题和四项研究内容描述')

# ============================================================
# SLIDE 10 (index 9): 模型架构 - 核心修改
# ============================================================
slide10 = prs.slides[9]

for shape in slide10.shapes:
    t = shape.text_frame.text if shape.has_text_frame else ''

    # 缺陷聚焦分支 - 加CBAM + 修正描述
    if '每个卷积块：Conv+IN+LeakyReLU' in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '每个卷积块：Conv+IN+LeakyReLU' in run.text:
                    run.text = run.text.replace(
                        '每个卷积块：Conv+IN+LeakyReLU',
                        '每个卷积块集成CBAM通道+空间双注意力'
                    )
                if 'Tanh激活输出缺陷图像' in run.text:
                    run.text = run.text.replace(
                        'Tanh激活输出缺陷图像',
                        'Tanh激活输出，全程AdaIN注入风格向量'
                    )
                if 'AdaIN注入风格向量' in run.text:
                    run.text = run.text.replace(
                        'AdaIN注入风格向量',
                        'AdaIN按层注入风格向量：粗层控形态、细层控纹理'
                    )
                if '控制缺陷类型/形态/严重程度' in run.text:
                    run.text = run.text.replace(
                        '控制缺陷类型/形态/严重程度',
                        '控制缺陷类型、位置、严重程度三要素'
                    )

    # 注意力融合模块 - 修正公式
    if '学习空间权重图' in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '学习空间权重图α∈[0,1]' in run.text:
                    run.text = run.text.replace(
                        '学习空间权重图α∈[0,1]',
                        '学习空间权重α，融合公式：α×缺陷 + (1-α)×背景'
                    )
                if '缺陷区域α→1，背景α→0' in run.text:
                    run.text = run.text.replace(
                        '缺陷区域α→1，背景α→0',
                        '经两层精炼卷积+Tanh消除拼接过渡痕迹'
                    )
                if '过渡区域平滑变化' in run.text:
                    run.text = run.text.replace(
                        '过渡区域平滑变化',
                        '缺陷与背景过渡区域平滑渐变'
                    )

    # 判别器描述修正
    if '关注全局结构一致性' in t and '尺度1' in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '关注全局结构一致性' in run.text:
                    run.text = run.text.replace(
                        '关注全局结构一致性',
                        '5层Conv+每层CBAM，全局结构一致性判别'
                    )

    if 'CBAM增强特征表达' in t and '尺度2' in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if 'CBAM增强特征表达' in run.text:
                    run.text = run.text.replace(
                        'CBAM增强特征表达',
                        '每层CBAM增强，平衡局部细节与全局结构'
                    )

    # 底部信息论表述修改
    if '双分支解耦设计从信息论角度' in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '双分支解耦设计从信息论角度' in run.text:
                    run.text = run.text.replace(
                        '双分支解耦设计从信息论角度显式约束缺陷与背景的条件独立性，有效降低学习任务的VC维，缓解模式崩溃隐患。多尺度判别器集成CBAM使判别器对细微缺陷响应更强。',
                        '双分支设计使缺陷生成与背景保持各司其职：缺陷分支从噪声学习纹理，背景分支编码输入结构，避免单一生成器同时处理两类异构信息。三尺度输出经1×1卷积融合得最终判别结果。'
                    )
                    break

print('Slide 10: 修正模型架构描述（CBAM/融合公式/判别器/底部表述）')

# ============================================================
# SLIDE 11 (index 10): 损失函数 - 修正感知损失描述 + 补充训练细节
# ============================================================
slide11 = prs.slides[10]
for shape in slide11.shapes:
    t = shape.text_frame.text if shape.has_text_frame else ''

    if 'L2距离' in t and 'VGG19' in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if 'L2距离' in run.text:
                    run.text = run.text.replace(
                        '提取VGG19 relu1_1~relu5_1层特征，计算生成与真实图像的L2距离。保持语义层面一致性。',
                        '提取VGG19的relu1_1~relu5_1共5层特征，计算MSE损失，从浅层纹理到高层语义多层约束。'
                    )
                    break

    if '基于深度特征的感知相似度计算' in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '基于深度特征的感知相似度计算' in run.text:
                    run.text = run.text.replace(
                        '基于深度特征的感知相似度计算，更符合人类视觉感知，保证生成图像视觉自然度。',
                        '基于VGG16的patch级感知距离（LPIPS），在深度特征空间衡量相似度，更符合人类视觉判断。'
                    )
                    break

    if '超参数通过Optuna框架自动优化' in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '超参数通过Optuna框架自动优化' in run.text:
                    run.text = run.text.replace(
                        '超参数通过Optuna框架自动优化（50次试验，以验证集FID最小化为目标）',
                        'Optuna贝叶斯优化（50次试验）：G_lr[1e-5,1e-3]、λ_gp[1,20]、λ_p[0.01,1.0]对数采样，目标最小化FID。训练：Epochs=100, Batch=8, n_critic=5, 256×256。'
                    )
                    break

print('Slide 11: 修正损失函数描述')

# ============================================================
# SLIDE 13 (index 12): 实验环境 - 修正框架和检测器
# ============================================================
slide13 = prs.slides[12]
for shape in slide13.shapes:
    t = shape.text_frame.text if shape.has_text_frame else ''

    if 'PyTorch 1.12.1' in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if 'PyTorch 1.12.1' in run.text:
                    run.text = run.text.replace('PyTorch 1.12.1', 'PyTorch 2.1')
                    break

    if 'PaDiM' in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if 'PaDiM' in run.text:
                    run.text = run.text.replace(
                        '检测器：PaDiM (ResNet18 + 马氏距离)',
                        '检测器：ResNet18/WRN50-2多尺度特征 + 马氏距离异常定位'
                    )
                    break

print('Slide 13: 修正框架版本和检测器描述')

# ============================================================
# SLIDE 14 (index 13): 生成质量 - 加PPS解释
# ============================================================
slide14 = prs.slides[13]
for shape in slide14.shapes:
    t = shape.text_frame.text if shape.has_text_frame else ''
    if 'PPS较StyleGAN2' in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if 'PPS较StyleGAN2' in run.text:
                    run.text = run.text.replace(
                        'PPS较StyleGAN2提升16.2%。',
                        'PPS（物理合理性得分）= 0.5×S_geo几何一致性 + 0.5×S_illum光照一致性，评估生成缺陷是否符合材料物理规律，较StyleGAN2提升16.2%。'
                    )
                    break

print('Slide 14: 添加PPS解释')

# ============================================================
# SLIDE 15 (index 14): 检测性能 - 补充消融细节
# ============================================================
slide15 = prs.slides[14]
for shape in slide15.shapes:
    t = shape.text_frame.text if shape.has_text_frame else ''
    if '双分支结构、各损失项均对性能有重要贡献' in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '双分支结构、各损失项均对性能有重要贡献' in run.text:
                    run.text = run.text.replace(
                        '双分支结构、各损失项均对性能有重要贡献。',
                        '去除CBAM后Pixel-AUC降幅最大，验证注意力机制对微缺陷感知的关键作用；去除感知损失导致纹理失真；各模块协同互补，完整模型最优。'
                    )
                    break

print('Slide 15: 补充消融实验细节')

# ============================================================
# SLIDE 16 (index 15): 系统前端 - 清理重复 + 修正多模式描述
# ============================================================
slide16 = prs.slides[15]
dup_removed = 0
for shape in slide16.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        # 删除底部重复的架构标签（位置靠下的重复元素）
        if shape.top > 7000000 and text in {
            '系统架构', '前端设计', 'Flask API', '增广引擎',
            'MVTec AD', 'HTML+CSS+JS Canvas', 'API 网关层', 'Focus-StyleGAN'
        }:
            sp = shape._element
            sp.getparent().remove(sp)
            dup_removed += 1

for shape in slide16.shapes:
    t = shape.text_frame.text if shape.has_text_frame else ''
    if 'GAN模型生成 + 数据集检索匹配' in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if 'GAN模型生成 + 数据集检索匹配' in run.text:
                    run.text = run.text.replace(
                        'GAN模型生成 + 数据集检索匹配  + 缺陷迁移堆叠，4张结果网格加8维评估面板',
                        '四种增广模式：GAN生成 | 真实缺陷迁移 | 检索匹配 | 缺陷堆叠。4张结果网格 + 8维评估面板'
                    )
                    break

    if '模式增广' in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '模式增广' in run.text:
                    run.text = run.text.replace('模式增广', '多模式增广引擎')
                    break

print(f'Slide 16: 删除{dup_removed}个重复元素，修正多模式描述')

# ============================================================
# SLIDE 18 (index 17): 结论 - 补充Web成果 + 修正不足
# ============================================================
slide18 = prs.slides[17]
for shape in slide18.shapes:
    t = shape.text_frame.text if shape.has_text_frame else ''

    if '完整增广评估体系' in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '完整增广评估体系' in run.text:
                    run.text = run.text.replace('完整增广评估体系', '完整评估与工程实现')
                if '图像质量(FID/IS/LPIPS/PPS)与下游检测性能(Pixel-AUC/PRO-AUC)双重验证' in run.text:
                    run.text = run.text.replace(
                        '图像质量(FID/IS/LPIPS/PPS)与下游检测性能(Pixel-AUC/PRO-AUC)双重验证',
                        '质量指标+检测性能双重评估 + Flask Web系统（4模式统一接入、Canvas后备、RESTful API）'
                    )
                    break

    if '计算资源需求较高' in t or '双分支需较多显存' in t:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if '计算资源需求较高' in run.text or '双分支需较多显存' in run.text:
                    run.text = run.text.replace(
                        '计算资源需求较高，双分支需较多显存',
                        '256×256分辨率下batch=8需约16GB显存（V100），消费级GPU部署需进一步优化'
                    )
                    break

print('Slide 18: 补充Web工程成果，修正不足描述')

# ============================================================
# SAVE
# ============================================================
prs.save(dst)
print(f'\n修改完成，已保存至: {dst}')
